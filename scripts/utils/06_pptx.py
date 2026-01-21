from pptx import Presentation
from pptx.util import Inches
from pptx.util import Cm
import pandas as pd
import os

list_file = os.path.join("pptx", "list.csv")
image_dir = os.path.join("pptx", "image")

data = pd.read_csv(list_file)

print(data)

####################
space_up = Cm(1)
space_bottom = Cm(1)
space_left = Cm(1)
space_right = Cm(1)

slide_width = Cm(33.87)
slide_height = Cm(19.05)

figure_width = Cm(20)
####################


prs = Presentation()

prs.slide_width = slide_width
prs.slide_height = slide_height

blank_slide_layout = prs.slide_layouts[6]

#figureの設定
figure_left, figure_top = space_left, space_up

#tableの設定
rows, cols = 9, 3
table_left, table_top = figure_left + figure_width, space_up
table_width, table_height = slide_width - (figure_left + space_right + figure_width), slide_height - (space_up + space_bottom)

for row in data.itertuples():
  slide = prs.slides.add_slide(blank_slide_layout)

  slide.shapes.add_picture(os.path.join(image_dir, f"{row.File}_plot_noScorpion.png"), figure_left, figure_top, width=figure_width)

  table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

  table.cell(0, 0).merge(table.cell(0, 2))
  table.cell(0, 0).text = row.File

  table.cell(1, 0).merge(table.cell(2, 0))
  table.cell(1, 0).text = f"{row.base_name}"
  table.cell(1, 1).text = "Chi2"
  table.cell(1, 2).text = f"{row.Chi2_base}"
  table.cell(2, 1).text = "DOF"
  table.cell(2, 2).text = f"{row.DOF_base}"

  table.cell(3, 0).merge(table.cell(4, 0))
  table.cell(3, 0).text = f"{row.comp_name}"
  table.cell(3, 1).text = "Chi2"
  table.cell(3, 2).text = f"{row.Chi2_comp}"
  table.cell(4, 1).text = "DOF"
  table.cell(4, 2).text = f"{row.DOF_comp}"

  table.cell(5, 0).merge(table.cell(5, 1))
  table.cell(5, 0).text = "Delta_Chi2"
  table.cell(5, 2).text = f"{row.Delta_Chi2}"

  table.cell(6, 0).merge(table.cell(6, 1))
  table.cell(6, 0).text = "f_val"
  table.cell(6, 2).text = f"{row.f_val}"

  table.cell(7, 0).merge(table.cell(7, 1))
  table.cell(7, 0).text = "p_val_ftest"
  table.cell(7, 2).text = f"{row.p_val_ftest}"

  table.cell(8, 0).merge(table.cell(8, 1))
  table.cell(8, 0).text = "p_val_mc"
  table.cell(8, 2).text = f"{row.p_val_mc}"

  for row in table.rows:
    row.height = Cm(1.0)
    row.width = Cm(3.0)

# 名前をつけて保存
prs.save(os.path.join("pptx", "Analysis_Report.pptx"))